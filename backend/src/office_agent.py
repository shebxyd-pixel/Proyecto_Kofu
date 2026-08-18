import os
import shutil
from typing import Any, Dict, List, Optional


from paths import BASE_DIR, ARCHIVOS_DIR, OUTPUT_DIR


def _ensure_dirs():
    os.makedirs(ARCHIVOS_DIR, exist_ok=True)
    os.makedirs(OUTPUT_DIR, exist_ok=True)


def _save_backup(src_path: str) -> str:
    _ensure_dirs()
    dest = os.path.join(ARCHIVOS_DIR, os.path.basename(src_path))
    shutil.copy2(src_path, dest)
    return os.path.abspath(dest)


class OfficeAgent:
    PPT_THEMES = {
        "professional": {"background": (240, 245, 250), "accent": (0, 102, 204), "text": (30, 30, 30)},
        "modern": {"background": (20, 20, 30), "accent": (0, 255, 200), "text": (240, 240, 240)},
        "vibrant": {"background": (255, 248, 220), "accent": (255, 87, 51), "text": (30, 30, 30)},
    }
    DOC_STYLES = {
        "professional": {"heading_font": "Arial", "body_font": "Calibri",
                         "heading_color": (0, 102, 204), "body_color": (30, 30, 30)},
        "modern": {"heading_font": "Segoe UI", "body_font": "Segoe UI",
                   "heading_color": (0, 200, 150), "body_color": (50, 50, 50)},
    }

    TEMPLATE_SCHEMAS = {
        "Coleccion Cuadratica": {
            "slides": [
                {"title": "{topic}", "layout": 0, "subtitle": "Presentación generada por Kofu", "background": True},
                {"title": "Introducción", "layout": 1, "text": "{intro}", "background": True},
                {"title": "Desarrollo", "layout": 1, "text": "{body}", "background": True},
                {"title": "Conclusiones", "layout": 5, "text": "{conclusion}", "background": True},
            ]
        },
        "Fibras Tejidas": {
            "slides": [
                {"title": "{topic}", "layout": 0, "subtitle": "Elaborado con Kofu", "background": True},
                {"title": "Contenido Principal", "layout": 1, "text": "{body}", "background": True},
                {"title": "Cierre", "layout": 5, "text": "{conclusion}", "background": True},
            ]
        },
        "default_word": {
            "sections": ["Introducción", "Desarrollo", "Análisis", "Conclusiones", "Referencias"]
        },
    }

    def __init__(self):
        self._win32_available = self._check_win32()

    @staticmethod
    def _check_win32() -> bool:
        try:
            import win32com.client  # noqa: F401
            return True
        except ImportError:
            return False

    def _convert_office_template(self, template_path: str, app_name: str, save_format: int,
                                  ext_from: str, ext_to: str) -> Optional[str]:
        import zipfile
        import tempfile
        import shutil
        template_path = os.path.abspath(template_path)
        temp_path = template_path.replace(ext_from, ext_to)
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                with zipfile.ZipFile(template_path, 'r') as zip_ref:
                    zip_ref.extractall(tmpdir)
                
                content_types_path = os.path.join(tmpdir, "[Content_Types].xml")
                with open(content_types_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                
                if ext_from == ".dotx":
                    content = content.replace(
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml",
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"
                    )
                elif ext_from == ".potx":
                    content = content.replace(
                        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
                    )
                    
                with open(content_types_path, 'w', encoding='utf-8') as f:
                    f.write(content)
                    
                with zipfile.ZipFile(temp_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
                    for root, dirs, files in os.walk(tmpdir):
                        for file in files:
                            file_path = os.path.join(root, file)
                            arcname = os.path.relpath(file_path, tmpdir)
                            zip_out.write(file_path, arcname)
            return temp_path
        except Exception:
            return None

    def get_template_schema(self, template_name: str) -> Optional[Dict]:
        if not template_name:
            return None
        base = os.path.splitext(template_name)[0]
        for key in self.TEMPLATE_SCHEMAS:
            if key.lower() in base.lower():
                return self.TEMPLATE_SCHEMAS[key]
        return None

    def create_powerpoint(self, output_path: str, slides_data: List[Dict[str, Any]],
                           template_path: Optional[str] = None, theme: str = "professional",
                           save_backup: bool = True) -> Dict[str, Any]:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        _ensure_dirs()
        temp_file = None
        if template_path and os.path.exists(template_path):
            path = template_path
            if path.endswith(".potx"):
                temp_file = self._convert_office_template(path, "PowerPoint.Application", 11, ".potx", "_temp.pptx")
                path = temp_file or None
            presentation = Presentation(path) if path else Presentation()
        else:
            presentation = Presentation()

        colors = self.PPT_THEMES.get(theme, self.PPT_THEMES["professional"])
        for slide_data in slides_data:
            layout = presentation.slide_layouts[slide_data.get("layout", 0)]
            slide = presentation.slides.add_slide(layout)

            if slide_data.get("background"):
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*colors["background"])

            for shape in slide.shapes:
                if hasattr(shape, "text") and "text" in slide_data:
                    shape.text = slide_data["text"]
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(18)
                            run.font.color.rgb = RGBColor(*colors["text"])

            if "title" in slide_data and slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
                for p in slide.shapes.title.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size, run.font.bold = Pt(36), True
                        run.font.color.rgb = RGBColor(*colors["accent"])

            if "subtitle" in slide_data and len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data["subtitle"]

        output_path = os.path.abspath(output_path)
        presentation.save(output_path)
        if temp_file and os.path.exists(temp_file):
            os.remove(temp_file)

        backup_path = None
        if save_backup:
            backup_path = _save_backup(output_path)

        return {"file_path": output_path, "backup_path": backup_path}

    def create_word_document(self, output_path: str, content_data: List[Dict[str, Any]],
                              template_path: Optional[str] = None, style: str = "professional",
                              save_backup: bool = True) -> Dict[str, Any]:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        _ensure_dirs()
        temp_file = None
        if template_path and os.path.exists(template_path):
            path = template_path
            if path.endswith(".dotx"):
                temp_file = self._convert_office_template(path, "Word.Application", 16, ".dotx", "_temp.docx")
                path = temp_file or None
            document = Document(path) if path else Document()
        else:
            document = Document()

        doc_style = self.DOC_STYLES.get(style, self.DOC_STYLES["professional"])

        import re

        for content in content_data:
            ctype = content["type"]
            if ctype == "heading":
                level = content.get("level", 1)
                heading = document.add_heading(content["text"], level=level)
                for run in heading.runs:
                    run.font.name = doc_style["heading_font"]
                    run.font.color.rgb = RGBColor(*doc_style["heading_color"])
                    run.font.bold = True
                    run.font.size = Pt({1: 24, 2: 18}.get(level, 14))
            elif ctype == "document_title":
                clean_title = content["text"].replace("**", "").replace("*", "").replace("#", "").strip()
                heading = document.add_heading(clean_title, level=0)
                for run in heading.runs:
                    run.font.name = doc_style["heading_font"]
                    run.font.color.rgb = RGBColor(*doc_style["heading_color"])
            elif ctype == "document_subtitle":
                clean_sub = content["text"].replace("**", "").replace("*", "").replace("#", "").strip()
                try:
                    p = document.add_paragraph(clean_sub, style='Subtitle')
                except Exception:
                    p = document.add_paragraph(clean_sub)
                    p.runs[0].font.bold = True
                for run in p.runs:
                    run.font.name = doc_style["heading_font"]
                    run.font.color.rgb = RGBColor(*doc_style["heading_color"])
            elif ctype == "paragraph":
                blocks = content["text"].split("\n")
                for block in blocks:
                    block = block.strip()
                    if not block:
                        continue
                    
                    if block.startswith("#"):
                        level = len(block) - len(block.lstrip("#"))
                        level = min(level, 9)
                        clean_text = block.lstrip("#").strip()
                        heading = document.add_heading(clean_text, level=level)
                        for run in heading.runs:
                            run.font.name = doc_style["heading_font"]
                            run.font.color.rgb = RGBColor(*doc_style["heading_color"])
                        continue

                    paragraph = document.add_paragraph()
                    paragraph.alignment = content.get("alignment", WD_ALIGN_PARAGRAPH.JUSTIFY)
                    
                    parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', block)
                    for part in parts:
                        if not part: continue
                        run = paragraph.add_run()
                        run.font.name, run.font.size = doc_style["body_font"], Pt(11)
                        run.font.color.rgb = RGBColor(*doc_style["body_color"])
                        
                        if part.startswith("**") and part.endswith("**"):
                            run.text = part[2:-2]
                            run.bold = True
                        elif part.startswith("*") and part.endswith("*"):
                            run.text = part[1:-1]
                            run.italic = True
                        else:
                            run.text = part
            elif ctype == "table":
                data = content["data"]
                table = document.add_table(rows=len(data), cols=len(data[0]))
                table.style = "Table Grid"
                for i, row in enumerate(data):
                    for j, cell_text in enumerate(row):
                        cell = table.rows[i].cells[j]
                        cell.text = cell_text
                        for p in cell.paragraphs:
                            for run in p.runs:
                                run.font.name = doc_style["body_font"]
                                if i == 0:
                                    run.font.bold = True
                                    run.font.color.rgb = RGBColor(*doc_style["heading_color"])
            elif ctype == "image" and os.path.exists(content["path"]):
                document.add_picture(content["path"], width=Inches(content.get("width", 6)))

        output_path = os.path.abspath(output_path)
        document.save(output_path)
        if temp_file and os.path.exists(temp_file):
            os.remove(temp_file)

        backup_path = None
        if save_backup:
            backup_path = _save_backup(output_path)

        return {"file_path": output_path, "backup_path": backup_path}
