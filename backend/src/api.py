"""
API unificada del asistente de IA local (Webtrest).
Fusiona: sanitización, corrección de texto, razonamiento (reglas + Ollama),
investigación web, procesamiento de archivos y generación de documentos Office
en un único backend ligero. Todas las dependencias pesadas u opcionales
(win32com, markitdown, bleach) se importan de forma perezosa y con
degradación controlada: si faltan, esa función se desactiva sin tumbar
el resto de la API.
"""

import io
import json
import os
import re
import urllib.parse
from typing import Any, Dict, List, Optional, Tuple
from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Permite peticiones desde cualquier origen (ideal para desarrollo local)
    allow_credentials=True,
    allow_methods=["*"],  # Permite todos los métodos (GET, POST, etc.)
    allow_headers=["*"],  # Permite todos los encabezados
)

import requests
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

DEFAULT_OLLAMA_MODEL = "gemma4:latest"
DEFAULT_OLLAMA_URL = "http://localhost:11434"


# ======================================================================
# 1. SANITIZACIÓN
# ======================================================================
class SanitizadorEntrada:
    """Limpia texto de patrones de inyección (script, SQL, comandos de shell)."""

    PATRONES_PELIGROSOS = [
        r"<script\b[^<]*(?:(?!<\/script>)<[^<]*)*<\/script>",
        r"javascript:", r"vbscript:", r"onload=", r"onclick=", r"onerror=",
        r"eval\s*\(", r"exec\s*\(", r"system\s*\(", r"shell\s*\(",
        r"drop\s+table", r"select\s+.*\s+from", r"insert\s+into", r"delete\s+from",
        r"rm\s+-rf", r"mkdir\s+", r"rmdir\s+", r"chmod\s+", r"sudo\s+",
        r"<iframe", r"<object", r"<embed", r"<svg", r"<link", r"<meta",
    ]
    ETIQUETAS_PERMITIDAS = ["b", "i", "u", "strong", "em", "p", "ul", "ol", "li", "h1", "h2", "h3"]

    def __init__(self):
        try:
            import bleach
            self._bleach = bleach
        except ImportError:
            self._bleach = None

    def limpiar_texto(self, texto: str) -> str:
        texto = str(texto)
        for patron in self.PATRONES_PELIGROSOS:
            texto = re.sub(patron, "", texto, flags=re.IGNORECASE)

        if self._bleach:
            texto = self._bleach.clean(
                texto, tags=self.ETIQUETAS_PERMITIDAS, attributes={}, strip=True, strip_comments=True
            )
        else:
            # Fallback sin bleach: elimina cualquier etiqueta HTML restante.
            texto = re.sub(r"<[^>]+>", "", texto)

        texto = re.sub(r"[\x00-\x1F\x7F-\x9F]", "", texto)
        return re.sub(r"\s+", " ", texto).strip()

    def es_entrada_segura(self, texto: str) -> bool:
        limpio = self.limpiar_texto(texto)
        return 0 < len(limpio) <= 10000

    def sanitizar_nombre_archivo(self, nombre: str) -> str:
        limpio = re.sub(r'[<>:"/\\|?*]', "", nombre)
        limpio = os.path.basename(limpio)
        return limpio[:255]


# ======================================================================
# 2. CORRECCIÓN DE ERRORES TIPOGRÁFICOS
# ======================================================================
class TypoCorrector:
    DICCIONARIO = {
        "powerpoint": ["powerpoin", "powerpint", "powepoint", "powerpiont", "ppt"],
        "word": ["wrod", "wor", "wrd", "wordd"],
        "documento": ["documeto", "document", "documnto", "doc"],
        "presentacion": ["presentacio", "presentacionn", "presntacion"],
        "plantilla": ["plantila", "plantllla", "plantil", "template"],
        "crear": ["crea", "crar", "hacer"],
        "investigar": ["investiga", "investgar", "buscar"],
        "consejos": ["consejo", "tips"],
        "ayuda": ["aydua", "help"],
    }
    FRASES_COMUNES = {
        "crear documento de word": ["crear doc de word", "hacer documento de word", "crear word"],
        "crear presentacion de powerpoint": ["crear ppt", "hacer presentacion powerpoint", "crear powerpoint"],
        "consejos de powerpoint": ["tips de powerpoint", "consejos ppt", "ayuda powerpoint"],
        "consejos de word": ["tips de word", "consejos doc", "ayuda word"],
        "investigar sobre": ["buscar sobre", "investiga sobre", "busca sobre"],
    }

    def correct_word(self, word: str) -> str:
        word_lower = word.lower()
        for correcta, typos in self.DICCIONARIO.items():
            if word_lower in typos or word_lower == correcta:
                return correcta
        return word

    def correct_text(self, text: str) -> str:
        text_lower = text.lower()
        for frase_correcta, variantes in self.FRASES_COMUNES.items():
            for variante in variantes:
                if variante in text_lower:
                    text_lower = text_lower.replace(variante, frase_correcta)
        palabras = re.findall(r"\w+", text_lower)
        return " ".join(self.correct_word(w) for w in palabras)


# ======================================================================
# 3. RAZONAMIENTO (reglas locales + Ollama como motor generativo)
# ======================================================================
class ReasoningStep:
    def __init__(self, step_num: int, thought: str, evidence: str = ""):
        self.step_num, self.thought, self.evidence = step_num, thought, evidence

    def __str__(self):
        return f"Paso {self.step_num}: {self.thought}" + (f" ({self.evidence})" if self.evidence else "")


class ChainOfThought:
    def __init__(self):
        self.steps: List[str] = []

    def add_step(self, step: str):
        self.steps.append(step)

    def get_chain(self) -> str:
        return "\n".join(f"{i + 1}. {s}" for i, s in enumerate(self.steps))

    def explain(self) -> str:
        return f"Proceso de razonamiento:\n{self.get_chain()}\n\nConclusión final derivada de los pasos anteriores."


class KnowledgeRule:
    def __init__(self, name: str, conditions: List[str], conclusion: str):
        self.name, self.conditions, self.conclusion = name, conditions, conclusion


class ReasoningEngine:
    """Motor basado en reglas: rápido, sin dependencias externas, siempre disponible."""

    def __init__(self):
        self.rules = [
            KnowledgeRule("crear_documento_pentesting",
                          ["necesita informacion sobre pentesting", "necesita un documento de seguridad"],
                          "crear_documento_pentesting"),
            KnowledgeRule("crear_presentacion_pentesting",
                          ["necesita una presentacion de pentesting", "presentacion de seguridad"],
                          "crear_presentacion_pentesting"),
            KnowledgeRule("investigar_tema",
                          ["necesita informacion sobre", "quiere saber sobre", "investiga", "busca"],
                          "investigar_tema"),
            KnowledgeRule("dar_consejos_office",
                          ["consejos de", "tips de", "trucos de", "como usar"],
                          "dar_consejos_office"),
            KnowledgeRule("explicar_pentesting", ["que es", "explica", "que son"], "explicar_pentesting"),
        ]
        self.facts: List[str] = []
        self.thinking_steps: List[ReasoningStep] = []

    def add_fact(self, fact: str):
        self.facts.append(fact.lower())

    def clear_facts(self):
        self.facts, self.thinking_steps = [], []

    def _extract_topic(self, text: str) -> str:
        for pattern in (r"sobre\s+(.+?)(?:\?|$|,|\.)", r"de\s+(.+?)(?:\?|$|,|\.)", r"acerca\s+de\s+(.+?)(?:\?|$|,|\.)"):
            match = re.search(pattern, text.lower())
            if match:
                return match.group(1).strip()
        return ""

    def _fallback_response(self, user_input: str) -> str:
        low = user_input.lower()
        if any(w in low for w in ("pentesting", "seguridad", "hacker")):
            return ("Sobre pentesting:\n• Metodologías: OWASP, PTES, OSSTMM, NIST\n"
                    "• Fases: Reconocimiento → Escaneo → Enumeración → Explotación → Post-explotación → Informes\n"
                    "• Herramientas: Nmap, Metasploit, Wireshark, Burp Suite\n\n¿Creo un documento o presentación?")
        if any(w in low for w in ("powerpoint", "presentación", "pptx", "ppt")):
            return "Puedo crear una presentación con temas professional, modern o vibrant. ¿Sobre qué tema?"
        if any(w in low for w in ("word", "documento", "docx")):
            return "Puedo crear un documento Word con estilo professional o modern. ¿Sobre qué tema?"
        return ("Puedo ayudarte con:\n• Pentesting\n• Documentos Word\n• Presentaciones PowerPoint\n"
                "• Investigación web\n• Consejos de Office\n\n¿Qué te gustaría hacer?")

    def reason(self, user_input: str) -> Tuple[str, List[ReasoningStep]]:
        self.clear_facts()
        self.add_fact(user_input)
        self.thinking_steps.append(ReasoningStep(1, "Analizando entrada del usuario", user_input))

        topic = self._extract_topic(user_input)
        if topic:
            self.thinking_steps.append(ReasoningStep(2, f"Tema identificado: {topic}"))

        for rule in self.rules:
            if any(cond in user_input.lower() for cond in rule.conditions):
                self.thinking_steps.append(ReasoningStep(3, f"Regla aplicada: {rule.name}"))
                return self._fallback_response(user_input), self.thinking_steps

        self.thinking_steps.append(ReasoningStep(3, "Ninguna regla específica aplicada"))
        return self._fallback_response(user_input), self.thinking_steps

    def get_thinking_process(self) -> str:
        return "\n".join(str(s) for s in self.thinking_steps)


class OllamaEngine:
    """Cliente robusto para un servidor Ollama local, con timeouts y errores explícitos."""

    def __init__(self, model_name: Optional[str] = None, base_url: Optional[str] = None, timeout: int = 5):
        self.base_url = (base_url or os.getenv("OLLAMA_BASE_URL", DEFAULT_OLLAMA_URL)).rstrip("/")
        self.model_name = model_name or os.getenv("OLLAMA_MODEL") or DEFAULT_OLLAMA_MODEL
        self.timeout = timeout
        self.available, self.available_models = self._check_availability()

    def _check_availability(self) -> Tuple[bool, List[str]]:
        try:
            resp = requests.get(f"{self.base_url}/api/tags", timeout=self.timeout)
            resp.raise_for_status()
            models = [m.get("name", "") for m in resp.json().get("models", [])]
            return True, models
        except Exception:
            return False, []

    def _generate(self, prompt: str, temperature: float = 0.7, num_predict: int = 1024, timeout: int = 120) -> str:
        if not self.available:
            raise RuntimeError("Ollama no está disponible.")
        payload = {
            "model": self.model_name, "prompt": prompt, "stream": False,
            "options": {"temperature": temperature, "num_predict": num_predict},
        }
        try:
            resp = requests.post(f"{self.base_url}/api/generate", json=payload, timeout=timeout)
            resp.raise_for_status()
            return resp.json()["response"].strip()
        except requests.exceptions.ConnectionError as e:
            raise RuntimeError("Ollama no responde. Verifica que el servicio esté corriendo.") from e
        except requests.exceptions.Timeout as e:
            raise RuntimeError(f"Tiempo de espera agotado tras {timeout}s.") from e
        except (KeyError, json.JSONDecodeError) as e:
            raise RuntimeError(f"Respuesta inesperada de Ollama: {e}") from e

    def razonar(self, prompt: str, system_prompt: Optional[str] = None) -> str:
        system_msg = system_prompt or (
            "Eres Webtrest, un asistente de IA para pentesting, documentos y presentaciones. "
            "Responde de forma clara y útil."
        )
        return self._generate(f"{system_msg}\n\nUSUARIO: {prompt}")

    def sanitizar_entrada(self, texto_sucio: str) -> Tuple[Dict[str, Any], bool]:
        fallback = {"texto_limpio": texto_sucio, "entidades_detectadas": [], "proceso_exitoso": False}
        if not self.available:
            return fallback, False
        prompt = (
            "Devuelve ÚNICAMENTE un JSON con los campos 'texto_limpio', 'entidades_detectadas' "
            f"y 'proceso_exitoso'.\n\nDATOS: '{texto_sucio}'"
        )
        try:
            raw = self._generate(prompt, temperature=0.1, num_predict=512, timeout=60)
            return json.loads(raw), True
        except (RuntimeError, json.JSONDecodeError):
            return fallback, False


def query_external_llm_fallback(prompt: str) -> Optional[str]:
    """Fallback opcional a un LLM en la nube (solo si OPENAI_API_KEY está configurada)."""
    api_key = os.getenv("OPENAI_API_KEY")
    if not api_key:
        return None
    try:
        resp = requests.post(
            "https://api.openai.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": "gpt-3.5-turbo", "messages": [{"role": "user", "content": prompt}],
                  "max_tokens": 500, "temperature": 0.7},
            timeout=30,
        )
        resp.raise_for_status()
        return resp.json()["choices"][0]["message"]["content"].strip()
    except Exception:
        return None


class HybridReasoningEngine(ReasoningEngine):
    """Reglas locales primero; si Ollama está disponible, lo usa para generar respuestas más ricas.
    Si Ollama falla, cae al fallback basado en reglas y, en último caso, a un LLM externo opcional."""

    def __init__(self, use_ollama: bool = True, model_name: Optional[str] = None):
        super().__init__()
        self.ollama_engine = OllamaEngine(model_name=model_name) if use_ollama else None
        self.chain_of_thought = ChainOfThought()

    def _respuesta_con_ollama(self, prompt: str, system_prompt: Optional[str] = None,
                               allow_external: bool = False) -> Optional[str]:
        """Ollama (local) se intenta siempre. El fallback a un LLM en la nube solo se
        usa si allow_external=True, es decir, si el modo activo es 'online'."""
        if self.ollama_engine and self.ollama_engine.available:
            try:
                return self.ollama_engine.razonar(prompt, system_prompt=system_prompt)
            except RuntimeError:
                pass
        return query_external_llm_fallback(prompt) if allow_external else None

    def direct_reason(self, user_input: str, allow_external: bool = False) -> Tuple[str, List[ReasoningStep]]:
        self.clear_facts()
        self.add_fact(user_input)
        self.thinking_steps.append(ReasoningStep(1, "Acceso directo sin sanitización", user_input))

        respuesta = self._respuesta_con_ollama(
            user_input,
            system_prompt="Eres un asistente de razonamiento directo. Responde de forma técnica y concisa.",
            allow_external=allow_external,
        )
        if respuesta:
            self.thinking_steps.append(ReasoningStep(2, "Respuesta generada con IA"))
            return respuesta, self.thinking_steps

        self.thinking_steps.append(ReasoningStep(2, "Respuesta directa sin motor generativo"))
        return user_input, self.thinking_steps

    def reason(self, user_input: str, allow_external: bool = False) -> Tuple[str, List[ReasoningStep]]:
        respuesta_reglas, pasos = super().reason(user_input)

        respuesta_ia = self._respuesta_con_ollama(user_input, allow_external=allow_external)
        if respuesta_ia:
            pasos.append(ReasoningStep(len(pasos) + 1, "Respuesta enriquecida con IA generativa"))
            return respuesta_ia, pasos

        return respuesta_reglas, pasos


# ======================================================================
# 4. INVESTIGACIÓN WEB
# ======================================================================
class WebResearcher:
    HEADERS = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) Chrome/120.0.0.0 Safari/537.36"}

    def __init__(self, google_api_key: Optional[str] = None, google_cx: Optional[str] = None):
        self.google_api_key = google_api_key or os.getenv("GOOGLE_API_KEY")
        self.google_cx = google_cx or os.getenv("GOOGLE_CX")

    def search_web(self, query: str, num_results: int = 5) -> List[Dict[str, str]]:
        for buscador in (self._search_google, self._search_duckduckgo, self._search_bing):
            try:
                results = buscador(query, num_results)
                if results:
                    return results
            except Exception:
                continue
        return []

    def _search_google(self, query: str, num_results: int) -> List[Dict[str, str]]:
        if not (self.google_api_key and self.google_cx):
            return []
        params = {"key": self.google_api_key, "cx": self.google_cx, "q": query, "num": min(num_results, 10)}
        resp = requests.get("https://www.googleapis.com/customsearch/v1", params=params, timeout=10)
        resp.raise_for_status()
        return [
            {"title": i.get("title", ""), "url": i.get("link", ""), "content": i.get("snippet", ""), "source": "Google"}
            for i in resp.json().get("items", [])
        ]

    def _search_duckduckgo(self, query: str, num_results: int) -> List[Dict[str, str]]:
        from bs4 import BeautifulSoup
        url = f"https://html.duckduckgo.com/html/?q={urllib.parse.quote(query)}"
        resp = requests.get(url, headers=self.HEADERS, timeout=10)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")
        results = []
        for div in soup.find_all("div", class_="result")[:num_results]:
            title_elem = div.find("a", class_="result__a")
            if not title_elem:
                continue
            snippet_elem = div.find("a", class_="result__snippet")
            results.append({
                "title": title_elem.get_text(strip=True),
                "url": title_elem.get("href", ""),
                "content": snippet_elem.get_text(strip=True) if snippet_elem else "",
                "source": "DuckDuckGo",
            })
        return results

    def _search_bing(self, query: str, num_results: int) -> List[Dict[str, str]]:
        from bs4 import BeautifulSoup
        url = f"https://www.bing.com/search?q={urllib.parse.quote(query)}"
        resp = requests.get(url, headers=self.HEADERS, timeout=10)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")
        results = []
        for item in soup.find_all("li", class_="b_algo")[:num_results]:
            title_elem, link_elem = item.find("h2"), item.find("a")
            if not (title_elem and link_elem):
                continue
            snippet_elem = item.find("p")
            href = link_elem.get("href", "")
            if href.startswith("http"):
                results.append({
                    "title": title_elem.get_text(strip=True), "url": href,
                    "content": snippet_elem.get_text(strip=True) if snippet_elem else "", "source": "Bing",
                })
        return results

    def generate_summary(self, topic: str, results: List[Dict[str, str]]) -> str:
        if not results:
            return f"No se encontró información sobre: {topic}"
        summary = f"📚 Resumen sobre: {topic}\n\n"
        for i, r in enumerate(results, 1):
            summary += f"🔹 Fuente {i} ({r.get('source', 'Desconocida')})\n   📝 {r['title']}\n"
            if r["content"]:
                summary += f"   ℹ️ {r['content'][:180]}...\n"
            summary += f"   🔗 {r['url']}\n\n"
        return summary.strip()


# ======================================================================
# 5. GENERACIÓN DE DOCUMENTOS OFFICE (PowerPoint + Word)
# ======================================================================
class OfficeAgent:
    """Genera .pptx y .docx. La conversión de plantillas .potx/.dotx requiere
    win32com (solo Windows + Office instalado); si no está disponible, se
    puede seguir generando documentos sin plantilla."""

    PPT_THEMES = {
        "professional": {"background": (240, 245, 250), "accent": (0, 102, 204), "text": (30, 30, 30)},
        "modern": {"background": (20, 20, 30), "accent": (0, 255, 200), "text": (240, 240, 240)},
        "vibrant": {"background": (255, 248, 220), "accent": (255, 87, 51), "text": (30, 30, 30)},
    }
    DOC_STYLES = {
        "professional": {"heading_font": "Arial", "body_font": "Calibri",
                          "heading_color": (0, 102, 204), "body_color": (30, 30, 30)},
        "modern": {"heading_font": "Segoe UI", "body_font": "Segoe UI",
                   "heading_color": (0, 200, 150), "body_color": (50, 50, 50)},
    }

    def __init__(self):
        self._win32_available = self._check_win32()

    @staticmethod
    def _check_win32() -> bool:
        try:
            import win32com.client  # noqa: F401
            return True
        except ImportError:
            return False

    def _convert_office_template(self, template_path: str, app_name: str, save_format: int, ext_from: str, ext_to: str) -> Optional[str]:
        """Convierte una plantilla (.potx/.dotx) a su formato editable usando Office vía COM."""
        if not self._win32_available:
            return None
        import win32com.client as win32
        template_path = os.path.abspath(template_path)
        temp_path = template_path.replace(ext_from, ext_to)
        app = None
        try:
            app = win32.Dispatch(app_name)
            app.Visible = False
            doc = app.Documents.Open(template_path) if app_name == "Word.Application" else app.Presentations.Open(template_path)
            doc.SaveAs(temp_path, save_format)
            doc.Close()
            return temp_path
        except Exception:
            return None
        finally:
            if app is not None:
                try:
                    app.Quit()
                except Exception:
                    pass

    # ---- PowerPoint ----
    def create_powerpoint(self, output_path: str, slides_data: List[Dict[str, Any]],
                           template_path: Optional[str] = None, theme: str = "professional") -> str:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        temp_file = None
        if template_path and os.path.exists(template_path):
            path = template_path
            if path.endswith(".potx"):
                temp_file = self._convert_office_template(path, "PowerPoint.Application", 11, ".potx", "_temp.pptx")
                path = temp_file or None
            presentation = Presentation(path) if path else Presentation()
        else:
            presentation = Presentation()

        colors = self.PPT_THEMES.get(theme, self.PPT_THEMES["professional"])
        for slide_data in slides_data:
            layout = presentation.slide_layouts[slide_data.get("layout", 0)]
            slide = presentation.slides.add_slide(layout)

            if slide_data.get("background"):
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*colors["background"])

            for shape in slide.shapes:
                if hasattr(shape, "text") and "text" in slide_data:
                    shape.text = slide_data["text"]
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(18)
                            run.font.color.rgb = RGBColor(*colors["text"])

            if "title" in slide_data and slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
                for p in slide.shapes.title.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size, run.font.bold = Pt(36), True
                        run.font.color.rgb = RGBColor(*colors["accent"])

            if "subtitle" in slide_data and len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data["subtitle"]

        output_path = os.path.abspath(output_path)
        presentation.save(output_path)
        if temp_file and os.path.exists(temp_file):
            os.remove(temp_file)
        return output_path

    # ---- Word ----
    def create_word_document(self, output_path: str, content_data: List[Dict[str, Any]],
                              template_path: Optional[str] = None, style: str = "professional") -> str:
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        temp_file = None
        if template_path and os.path.exists(template_path):
            path = template_path
            if path.endswith(".dotx"):
                temp_file = self._convert_office_template(path, "Word.Application", 16, ".dotx", "_temp.docx")
                path = temp_file or None
            document = Document(path) if path else Document()
        else:
            document = Document()

        doc_style = self.DOC_STYLES.get(style, self.DOC_STYLES["professional"])

        for content in content_data:
            ctype = content["type"]
            if ctype == "heading":
                level = content.get("level", 1)
                heading = document.add_heading(content["text"], level=level)
                for run in heading.runs:
                    run.font.name = doc_style["heading_font"]
                    run.font.color.rgb = RGBColor(*doc_style["heading_color"])
                    run.font.bold = True
                    run.font.size = Pt({1: 24, 2: 18}.get(level, 14))
            elif ctype == "paragraph":
                paragraph = document.add_paragraph()
                paragraph.alignment = content.get("alignment", WD_ALIGN_PARAGRAPH.JUSTIFY)
                run = paragraph.add_run(content["text"])
                run.font.name, run.font.size = doc_style["body_font"], Pt(11)
                run.font.color.rgb = RGBColor(*doc_style["body_color"])
            elif ctype == "table":
                data = content["data"]
                table = document.add_table(rows=len(data), cols=len(data[0]))
                table.style = "Table Grid"
                for i, row in enumerate(data):
                    for j, cell_text in enumerate(row):
                        cell = table.rows[i].cells[j]
                        cell.text = cell_text
                        for p in cell.paragraphs:
                            for run in p.runs:
                                run.font.name = doc_style["body_font"]
                                if i == 0:
                                    run.font.bold = True
                                    run.font.color.rgb = RGBColor(*doc_style["heading_color"])
            elif ctype == "image" and os.path.exists(content["path"]):
                document.add_picture(content["path"], width=Inches(content.get("width", 6)))

        output_path = os.path.abspath(output_path)
        document.save(output_path)
        if temp_file and os.path.exists(temp_file):
            os.remove(temp_file)
        return output_path


# ======================================================================
# 6. PROCESAMIENTO DE ARCHIVOS (MarkItDown, opcional)
# ======================================================================
class FileProcessor:
    SUPPORTED_EXTENSIONS = {
        ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv", ".txt", ".md",
        ".html", ".htm", ".xml", ".json", ".yaml", ".yml", ".rtf",
        ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff",
        ".wav", ".mp3", ".m4a", ".ogg", ".flac", ".epub", ".zip", ".msg", ".ipynb",
    }

    def __init__(self):
        try:
            from markitdown import MarkItDown
            self._markitdown = MarkItDown(enable_plugins=False)
        except ImportError:
            self._markitdown = None

    @property
    def available(self) -> bool:
        return self._markitdown is not None

    def process_local_file(self, file_path: str) -> Dict[str, Any]:
        if not self.available:
            return {"success": False, "error": "MarkItDown no está instalado.", "filename": os.path.basename(file_path)}
        try:
            result = self._markitdown.convert_local(file_path)
            return {
                "success": True,
                "text_content": result.text_content,
                "filename": os.path.basename(file_path),
                "file_type": os.path.splitext(file_path)[1].lstrip(".") or "unknown",
            }
        except Exception as e:
            return {"success": False, "error": str(e), "filename": os.path.basename(file_path)}

    @classmethod
    def is_supported_file(cls, filename: str) -> bool:
        return os.path.splitext(filename)[1].lower() in cls.SUPPORTED_EXTENSIONS


# ======================================================================
# 7. ORQUESTADOR PRINCIPAL
# ======================================================================
class KnowledgeBase:
    POWERPOINT_TIPS = [
        "Usa la regla de 6x6: máximo 6 líneas por diapositiva y 6 palabras por línea",
        "Aplica temas consistentes para mantener el diseño profesional",
        "Usa transiciones moderadas para no distraer",
        "Las animaciones deben ser funcionales, no decorativas",
        "Incluye notas del orador para información adicional",
    ]
    WORD_TIPS = [
        "Usa estilos predefinidos para títulos y párrafos",
        "Aplica sangría francesa para listas de referencias",
        "Usa la función de navegación para moverte rápidamente",
        "Habilita la revisión de cambios para colaboraciones",
        "Usa plantillas (.dotx) para mantener la coherencia",
    ]
    PENTESTING = {
        "metodologias": ["OWASP", "PTES", "OSSTMM", "NIST"],
        "herramientas": ["Nmap", "Metasploit", "Wireshark", "Burp Suite", "Hydra", "John the Ripper"],
        "fases": ["Reconocimiento", "Escaneo", "Enumeración", "Explotación", "Post-explotación", "Informes"],
    }


class AIAssistant:
    def __init__(self, use_ollama: bool = True, ollama_model: str = DEFAULT_OLLAMA_MODEL):
        self.sanitizador = SanitizadorEntrada()
        self.typo_corrector = TypoCorrector()
        self.web_researcher = WebResearcher()
        self.office_agent = OfficeAgent()
        self.file_processor = FileProcessor()
        self.knowledge_base = KnowledgeBase()
        self.reasoning_engine = HybridReasoningEngine(use_ollama=use_ollama, model_name=ollama_model)
        self.chain_of_thought = ChainOfThought()

    @staticmethod
    def _permite_internet(modo: str) -> bool:
        """'local' = 100% offline (solo Ollama + reglas). 'online' = permite investigación
        web y fallback a un LLM en la nube además de Ollama."""
        return modo == "online"

    def _preparar_entrada(self, user_input: str) -> str:
        """Sanitiza contra inyección y corrige errores tipográficos por diccionario.
        Se ejecuta SIEMPRE antes de que el texto llegue a Ollama o al motor de reglas,
        tanto en modo local como online."""
        if not self.sanitizador.es_entrada_segura(user_input):
            user_input = self.sanitizador.limpiar_texto(user_input)
        return self.typo_corrector.correct_text(user_input)

    def process_request(self, user_input: str, modo: str = "local") -> Tuple[str, List[ReasoningStep]]:
        entrada = self._preparar_entrada(user_input)
        self.chain_of_thought.steps = []
        self.chain_of_thought.add_step(f"Analizando solicitud ({modo}): {entrada}")
        response, steps = self.reasoning_engine.reason(entrada, allow_external=self._permite_internet(modo))
        self.chain_of_thought.add_step("Generando respuesta final")
        return response, steps

    def direct_reason(self, user_input: str, modo: str = "local") -> Tuple[str, List[ReasoningStep]]:
        entrada = self._preparar_entrada(user_input)
        self.chain_of_thought.steps = []
        self.chain_of_thought.add_step(f"Acceso directo ({modo}): {entrada}")
        response, steps = self.reasoning_engine.direct_reason(entrada, allow_external=self._permite_internet(modo))
        self.chain_of_thought.add_step("Respuesta directa generada")
        return response, steps

    def research_topic(self, topic: str, modo: str = "online") -> str:
        if not self._permite_internet(modo):
            return "La investigación web requiere modo 'online' (necesita conexión a internet)."
        results = self.web_researcher.search_web(topic)
        return self.web_researcher.generate_summary(topic, results)

    def get_office_tips(self, software: str) -> List[str]:
        return {"powerpoint": self.knowledge_base.POWERPOINT_TIPS,
                "word": self.knowledge_base.WORD_TIPS}.get(software.lower(), [])

    def create_presentation(self, topic: str, output_path: str, theme: str = "professional", modo: str = "online") -> str:
        research = self.research_topic(topic, modo)
        slides = [
            {"title": f"Presentación sobre: {topic}", "layout": 0, "subtitle": "Generado automáticamente", "background": True},
            {"title": "Introducción", "layout": 1, "text": research[:300] + "...", "background": True},
            {"title": "Conclusiones", "layout": 5, "text": "Para más información, consulte las fuentes originales.", "background": True},
        ]
        return self.office_agent.create_powerpoint(output_path, slides, theme=theme)

    def create_document(self, topic: str, output_path: str, style: str = "professional", modo: str = "online") -> str:
        research = self.research_topic(topic, modo)
        content = [
            {"type": "heading", "text": f"Documento sobre: {topic}", "level": 1},
            {"type": "paragraph", "text": research},
            {"type": "heading", "text": "Conocimiento adicional", "level": 2},
            {"type": "paragraph", "text": "Este documento fue generado automáticamente con información de fuentes públicas."},
        ]
        return self.office_agent.create_word_document(output_path, content, style=style)

    def digest_file(self, file_path: str, output_path: Optional[str] = None,
                     instrucciones: Optional[str] = None) -> Dict[str, Any]:
        """Pipeline completo del modo local para 'digerir' un archivo:
        1. MarkItDown traduce el archivo (pdf, docx, imagen, audio, etc.) a texto/Markdown.
        2. El texto se sanitiza y se corrige (typo_corrector) ANTES de tocar el modelo.
        3. Ollama procesa/estructura ese texto ya limpio.
        4. El resultado se guarda en output_path si se especifica.
        """
        extraido = self.file_processor.process_local_file(file_path)
        if not extraido.get("success"):
            return extraido

        texto_preparado = self._preparar_entrada(extraido["text_content"])

        prompt = instrucciones or "Resume y estructura el siguiente contenido de forma clara y organizada:"
        digestion = texto_preparado
        if self.reasoning_engine.ollama_engine and self.reasoning_engine.ollama_engine.available:
            try:
                digestion = self.reasoning_engine.ollama_engine.razonar(f"{prompt}\n\n{texto_preparado}")
            except RuntimeError:
                pass  # sin Ollama disponible, se conserva el texto ya sanitizado/corregido

        resultado = {"success": True, "filename": extraido["filename"], "digestion": digestion}

        if output_path:
            os.makedirs(os.path.dirname(os.path.abspath(output_path)) or ".", exist_ok=True)
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(digestion)
            resultado["output_path"] = os.path.abspath(output_path)

        return resultado


# ======================================================================
# 8. API HTTP (FastAPI)
# ======================================================================
app = FastAPI(title="Webtrest API", version="2.0.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

assistant = AIAssistant()


class ChatRequest(BaseModel):
    message: str
    direct: bool = False
    modo: str = "local"  # "local" (solo Ollama, sin internet) | "online" (permite web + fallback en la nube)


class TopicRequest(BaseModel):
    topic: str
    modo: str = "online"  # investigar requiere internet


class PresentationRequest(BaseModel):
    topic: str
    output_path: str
    theme: str = "professional"
    modo: str = "online"


class DocumentRequest(BaseModel):
    topic: str
    output_path: str
    style: str = "professional"
    modo: str = "online"


class TipsRequest(BaseModel):
    software: str


class DigestRequest(BaseModel):
    file_path: str
    output_path: Optional[str] = None
    instrucciones: Optional[str] = None


class CorrectRequest(BaseModel):
    text: str


@app.get("/health")
def health():
    return {
        "status": "ok",
        "ollama_available": assistant.reasoning_engine.ollama_engine.available if assistant.reasoning_engine.ollama_engine else False,
        "ollama_model": assistant.reasoning_engine.ollama_engine.model_name if assistant.reasoning_engine.ollama_engine else None,
        "office_templates_available": assistant.office_agent._win32_available,
        "file_processor_available": assistant.file_processor.available,
        "modos_disponibles": ["local", "online"],
    }


@app.post("/chat")
def chat(req: ChatRequest):
    try:
        response, steps = (
            assistant.direct_reason(req.message, modo=req.modo) if req.direct
            else assistant.process_request(req.message, modo=req.modo)
        )
        return {"response": response, "steps": [str(s) for s in steps], "modo": req.modo}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/research")
def research(req: TopicRequest):
    try:
        return {"summary": assistant.research_topic(req.topic, modo=req.modo)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/office/powerpoint")
def create_powerpoint(req: PresentationRequest):
    try:
        path = assistant.create_presentation(req.topic, req.output_path, req.theme, modo=req.modo)
        return {"path": path}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/office/word")
def create_word(req: DocumentRequest):
    try:
        path = assistant.create_document(req.topic, req.output_path, req.style, modo=req.modo)
        return {"path": path}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/office/tips")
def office_tips(req: TipsRequest):
    return {"tips": assistant.get_office_tips(req.software)}


@app.post("/files/digest")
def digest_file(req: DigestRequest):
    """Pipeline del modo local: MarkItDown (archivo -> texto) -> sanitización/typo_corrector
    -> Ollama (digestión) -> archivo de salida opcional."""
    if not assistant.file_processor.available:
        raise HTTPException(status_code=503, detail="MarkItDown no está instalado. Ejecuta: pip install markitdown")
    resultado = assistant.digest_file(req.file_path, req.output_path, req.instrucciones)
    if not resultado.get("success"):
        raise HTTPException(status_code=422, detail=resultado.get("error", "No se pudo procesar el archivo."))
    return resultado


@app.post("/text/correct")
def correct_text(req: CorrectRequest):
    """Expone la sanitización + corrección tipográfica que corre antes de cada solicitud,
    útil para depurar por qué el modelo recibió cierto texto."""
    sanitizado = req.text if assistant.sanitizador.es_entrada_segura(req.text) else assistant.sanitizador.limpiar_texto(req.text)
    corregido = assistant.typo_corrector.correct_text(sanitizado)
    return {"original": req.text, "sanitizado": sanitizado, "corregido": corregido}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
